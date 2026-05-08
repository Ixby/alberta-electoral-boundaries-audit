from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# SLIDE 1 — ORIENTATION
slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
slide.shapes.title.text = "Phase 1 Complete"
tf = slide.placeholders[1].text_frame
tf.text = "• Official Elections Alberta shapefiles received; analysis complete"
tf.add_paragraph().text = "• 50,000 neutral maps generated, all 3 pre-registered test channels executed"
tf.add_paragraph().text = "• 4 OSF pre-registrations filed; Fisher combined result in hand"
tf.add_paragraph().text = "\nA computational comparison of two proposed 2026 maps against a distribution of 50,000 maps drawn by purely population-constrained random redistricting."

# SLIDE 2 — THE QUESTION
slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
slide.shapes.title.text = "One Question"
slide.placeholders[1].text = "Does the minority map fall within the distribution of maps a neutral, unbiased redistricting process would produce?\n\n(The majority map and the minority map start from the same legal constraints. The question is whether their outcomes are statistically interchangeable.)"

# SLIDE 3 — METHOD
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "50,000 Neutral Maps"
tf = slide.placeholders[1].text_frame
tf.text = "• ReCom algorithm (MGGG/MIT) — random redistricting, no partisan knowledge"
tf.add_paragraph().text = "• 4,765 VA polygons (Elections Alberta official geometry)"
tf.add_paragraph().text = "• ±25% population constraint — same rule the commission used"
tf.add_paragraph().text = "• Scored on 4 metrics: Efficiency Gap, Mean-Median, Declination, Seats@50/50"
tf.add_paragraph().text = "• Pre-registered before shapefiles arrived; thresholds locked to drand beacon\n"
p = tf.add_paragraph()
p.text = "Pre-registration = the goalposts cannot move after seeing the data"
p.font.bold = True

# SLIDE 4 — CHANNEL 1: PARTISAN JOINT TAIL
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Channel 1: Joint Partisan Tail"
tf = slide.placeholders[1].text_frame
tf.text = "Minority 2026: Mahalanobis = 6.11 | p = 1.60 × 10⁻⁷ | Outside null"
tf.add_paragraph().text = "Majority 2026: Mahalanobis = 2.69 | p = 0.125 | Within null"
tf.add_paragraph().text = "2019 Enacted: Mahalanobis = 3.56 | p = 0.013 | Modest lean"
tf.add_paragraph().text = "\nEnsemble centre is already +1.6 pp UCP-leaning due to Alberta's geography. The minority map is an outlier beyond that natural baseline."
p = tf.add_paragraph()
p.text = "\np = 1.60 × 10⁻⁷"
p.font.size = Pt(40)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# SLIDE 5 — CHANNEL 2 + FISHER
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Channel 2 + Fisher Combined"
tf = slide.placeholders[1].text_frame
tf.text = "• SZAT: 2,108 of 4,765 VAs assigned differently between the two maps"
tf.add_paragraph().text = "• Bootstrap null (10,000 reshuffles): p = 0.0044"
tf.add_paragraph().text = "• Fisher combined (Channels 1 + 2): p = 1.55 × 10⁻⁸"
p = tf.add_paragraph()
p.text = "\n1 in 64 million\n"
p.font.size = Pt(40)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
tf.add_paragraph().text = "Fisher combined was the pre-registered primary outcome — not post-hoc"
tf.add_paragraph().text = "Directional context: Minority map population variance: 98.9th percentile of ensemble"

# SLIDE 6 — CHANNEL 3 (NULL)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Channel 3: Neighbour-Drain — Null Result"
tf = slide.placeholders[1].text_frame
tf.text = "• Tests pack-crack coupling mechanism (not outcome)"
tf.add_paragraph().text = "• Minority map: within null (p = 0.134)"
tf.add_paragraph().text = "• Majority map: anomalously clean (p < 0.0001, inverted)"
tf.add_paragraph().text = "• Pre-registered predictions A and B not confirmed\n"
p = tf.add_paragraph()
p.text = "The mechanism is not confirmed. The outcome findings (Ch1, Ch2) stand independently. All three channels are disclosed."
p.font.bold = True

# SLIDE 7 — JUSTIFICATION TESTS
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Five Justification Tests — All Fail"
tf = slide.placeholders[1].text_frame
tf.text = "Could the minority map's configurations be explained by population math or area law?\n"
tf.add_paragraph().text = "1. Olds–Three Hills–Didsbury"
tf.add_paragraph().text = "2. Rocky Mountain House–Banff extension"
tf.add_paragraph().text = "3. Airdrie 4-way split"
tf.add_paragraph().text = "4. Red Deer 4 districts"
tf.add_paragraph().text = "5. Chestermere into Calgary-Peigan\n"
p = tf.add_paragraph()
p.text = "The math does not force these choices. Cleaner alternatives exist in every case, and the majority map found them."
p.font.bold = True

# SLIDE 8 — SCOPE LIMITS
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What This Does and Does Not Say"
tf = slide.placeholders[1].text_frame
tf.text = "DOES SAY:"
tf.add_paragraph().text = "• The minority map sits at an extreme statistical position relative to neutral redistricting"
tf.add_paragraph().text = "• Two independent pre-registered tests reject the null"
tf.add_paragraph().text = "• Five contested configurations have simpler neutral alternatives\n"
tf.add_paragraph().text = "DOES NOT SAY:"
tf.add_paragraph().text = "• Who drew the map with what intent"
tf.add_paragraph().text = "• Whether this constitutes a gerrymander in law"
tf.add_paragraph().text = "• Anything about individual MLAs or cabinet ministers\n"
p = tf.add_paragraph()
p.text = "This is a statistical instrument. Legal and political conclusions belong to others."
p.font.italic = True

# SLIDE 9 — PHASE 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Phase 2: Pre-Registered and Waiting"
tf = slide.placeholders[1].text_frame
tf.text = "• Lunty committee (91-seat map, tabling late 2026) already pre-registered"
tf.add_paragraph().text = "• AsPredicted #289,455 / OSF osf.io/qsgy8"
tf.add_paragraph().text = "• 17-point signature scorecard locked to April 27 drand beacon"
tf.add_paragraph().text = "• Within 72 hours of tabling: automatic evaluation\n"
p = tf.add_paragraph()
p.text = "• The criteria exist before the map exists\n"
p.font.bold = True
tf.add_paragraph().text = "• Total chain: 4 pre-registrations, all public since May 7, 2026"

# SLIDE 10 — THE ASK
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Ask"
tf = slide.placeholders[1].text_frame
tf.text = "1. Academic review of methodology before SSRN submission"
tf.add_paragraph().text = "2. Brief commentary or framing note for the public report"
tf.add_paragraph().text = "3. Referral into CPSA or Canadian electoral studies community\n\n"
p = tf.add_paragraph()
p.text = "The Phase 1 analysis is done. The Phase 2 trigger is pre-registered. I'm not going to the press until I've had academic review."
p.font.bold = True

prs.save('presentation_dr_bratt.pptx')
